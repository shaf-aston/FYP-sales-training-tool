"""
Comprehensive Document Extractor for Final Year Project
========================================================
Extracts text from ALL document types in Aston Docs folder:
- PDF files (.pdf)
- Word documents (.docx, .dotx)
- PowerPoint presentations (.pptx)
- Existing text files (.txt) - consolidates them

Output: Single consolidated file for AI/LLM consumption
"""
import os
import hashlib
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path

# Try importing optional libraries
try:
    import fitz  # PyMuPDF for PDFs
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False
    print("⚠️ PyMuPDF not installed - PDF extraction disabled")

try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False
    print("⚠️ python-pptx not installed - PPTX extraction disabled")

# Configuration
BASE_FOLDER = Path(r"c:\Users\Shaf\Downloads\Final Year Project pack folder\Sales roleplay chatbot\Aston Docs")
OUTPUT_FILE = BASE_FOLDER / "CONSOLIDATED_ALL_DOCS.txt"
WORD_NAMESPACE = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'

# Track all extracted content
all_content = []

def add_section(title, content, source_file=""):
    """Add a section to the consolidated output."""
    separator = "=" * 80
    header = f"\n{separator}\n📄 {title}\n"
    if source_file:
        header += f"📁 Source: {source_file}\n"
    header += f"{separator}\n\n"
    all_content.append(header + content + "\n")

def extract_pdf(file_path):
    """Extract text from PDF using PyMuPDF."""
    if not PDF_SUPPORT:
        return f"[PDF extraction not available - install pymupdf]"
    
    try:
        doc = fitz.open(file_path)
        text_parts = []
        
        for page_num, page in enumerate(doc, 1):
            text = page.get_text()
            if text.strip():
                text_parts.append(f"--- Page {page_num} ---\n{text}")
        
        doc.close()
        return "\n".join(text_parts)
    except Exception as e:
        return f"ERROR: {str(e)}"

def extract_docx(file_path):
    """Extract text from DOCX/DOTX using ZIP/XML extraction."""
    try:
        full_text = []
        
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            all_files = zip_ref.namelist()
            
            # Extract from document.xml
            for xml_file in ['word/document.xml', 'word/document2.xml']:
                if xml_file in all_files:
                    content = zip_ref.read(xml_file)
                    root = ET.fromstring(content)
                    
                    # Extract paragraphs in order
                    for para in root.iter(f'{WORD_NAMESPACE}p'):
                        para_text = []
                        for t_elem in para.iter(f'{WORD_NAMESPACE}t'):
                            if t_elem.text:
                                para_text.append(t_elem.text)
                        if para_text:
                            full_text.append(''.join(para_text))
                    break
        
        return "\n".join(full_text)
    except Exception as e:
        return f"ERROR: {str(e)}"

def extract_pptx(file_path):
    """Extract text from PowerPoint presentations."""
    if not PPTX_SUPPORT:
        return "[PPTX extraction not available - install python-pptx]"
    
    try:
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"--- Slide {slide_num} ---"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if len(slide_text) > 1:  # More than just the header
                text_parts.append("\n".join(slide_text))
        
        return "\n\n".join(text_parts)
    except Exception as e:
        return f"ERROR: {str(e)}"

def extract_txt(file_path):
    """Read existing text files."""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        return f"ERROR: {str(e)}"

def get_file_hash(file_path):
    """Get MD5 hash for duplicate detection."""
    hash_md5 = hashlib.md5()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(4096), b""):
            hash_md5.update(chunk)
    return hash_md5.hexdigest()

def process_file(file_path):
    """Process a single file based on its extension."""
    ext = file_path.suffix.lower()
    name = file_path.name
    
    if ext == '.pdf':
        return name, extract_pdf(file_path)
    elif ext in ['.docx', '.dotx']:
        return name, extract_docx(file_path)
    elif ext == '.pptx':
        return name, extract_pptx(file_path)
    elif ext == '.txt':
        # Skip our own output files
        if 'CONSOLIDATED' in name or 'extracted' in name.lower():
            return None, None
        return name, extract_txt(file_path)
    else:
        return None, None

def scan_directory(folder, category=""):
    """Recursively scan a directory for documents."""
    folder = Path(folder)
    
    for item in sorted(folder.iterdir()):
        if item.is_dir():
            # Skip certain directories
            if item.name in ['__pycache__', 'text ppts', 'extracted_content']:
                continue
            scan_directory(item, category=item.name)
        elif item.is_file():
            name, content = process_file(item)
            if name and content:
                title = f"{category}/{name}" if category else name
                add_section(title, content, str(item))
                print(f"✓ Extracted: {title}")

def main():
    print("=" * 60)
    print("COMPREHENSIVE DOCUMENT EXTRACTOR")
    print("Final Year Project - Sales Roleplay Chatbot")
    print("=" * 60 + "\n")
    
    # Add header to consolidated output
    all_content.append("""
################################################################################
#                                                                              #
#              CONSOLIDATED DOCUMENTATION - SALES ROLEPLAY CHATBOT             #
#                          Final Year Project                                  #
#                                                                              #
#  This file contains ALL extracted text from project documentation.           #
#  Generated for AI/LLM consumption to assist with report writing.             #
#                                                                              #
################################################################################

TABLE OF CONTENTS:
1. Student Written Documents (Analyses.docx, FINAL YEAR PROJECT.docx)
2. Mark Scheme & Guidelines (CS3IP Mark Scheme, Report Templates)
3. Software Project Management Concepts (Week 1-11 lectures)
4. Academic Writing Guidelines

""")
    
    # Process main Aston Docs folder
    print("📂 Scanning Aston Docs folder...")
    scan_directory(BASE_FOLDER)
    
    # Process Project management subfolder
    pm_folder = BASE_FOLDER / "Project management"
    if pm_folder.exists():
        print("\n📂 Scanning Project Management folder...")
        scan_directory(pm_folder, category="Project Management")
    
    # Write consolidated output
    print(f"\n📝 Writing consolidated output to: {OUTPUT_FILE}")
    with open(OUTPUT_FILE, 'w', encoding='utf-8') as f:
        f.write("\n".join(all_content))
    
    # Summary
    print("\n" + "=" * 60)
    print("EXTRACTION COMPLETE")
    print("=" * 60)
    print(f"Total sections extracted: {len(all_content) - 1}")  # -1 for header
    print(f"Output file: {OUTPUT_FILE}")
    print(f"Output size: {OUTPUT_FILE.stat().st_size / 1024:.1f} KB")
    
    # Print document categories found
    print("\n📋 Documents processed by type:")
    extensions = {}
    for item in BASE_FOLDER.rglob('*'):
        if item.is_file():
            ext = item.suffix.lower()
            extensions[ext] = extensions.get(ext, 0) + 1
    
    for ext, count in sorted(extensions.items()):
        print(f"   {ext}: {count} files")

if __name__ == "__main__":
    main()
