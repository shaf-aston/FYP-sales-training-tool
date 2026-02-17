"""Extractor for Aston Docs — PDF & PPTX with optional OCR.

This is a compact, clean implementation that scans a directory recursively,
extracts text and OCR content, and writes outputs to `extracted_content`.
"""

from pathlib import Path
import argparse
import io
import logging
import zipfile
import xml.etree.ElementTree as ET
from datetime import datetime
from typing import Optional, Tuple

try:
    import fitz
except Exception:
    fitz = None

try:
    from PIL import Image
except Exception:
    Image = None

try:
    import pytesseract
except Exception:
    pytesseract = None

LOG = logging.getLogger(__name__)


def setup_output_directory(base: Path) -> Path:
    out = base / "extracted_content"
    out.mkdir(parents=True, exist_ok=True)
    return out


def extract_text_from_pdf(pdf_path: Path) -> str:
    if fitz is None:
        return ""
    try:
        doc = fitz.open(pdf_path)
        parts = []
        for i, page in enumerate(doc, 1):
            t = page.get_text() or ""
            if t.strip():
                parts.append(f"\n{'='*80}\nPage {i}\n{'='*80}\n" + t)
        doc.close()
        return "\n".join(parts)
    except Exception:
        return ""


def extract_images_from_pdf(pdf_path: Path, ocr_lang: Optional[str] = None, ocr_config: Optional[str] = None) -> Tuple[str, int]:
    if fitz is None or Image is None:
        return "[Image extraction not available]", 0
    try:
        doc = fitz.open(pdf_path)
        blocks = []
        count = 0
        for pno, page in enumerate(doc, 1):
            imgs = page.get_images(full=True)
            for idx, img in enumerate(imgs, 1):
                try:
                    xref = img[0]
                    base = doc.extract_image(xref)
                    img_bytes = base.get('image')
                    ext = base.get('ext', 'png')
                    if not img_bytes:
                        continue
                    image = Image.open(io.BytesIO(img_bytes)).convert('RGB')
                    count += 1
                    w, h = image.size
                    try:
                        ocr_text = pytesseract.image_to_string(image, lang=ocr_lang, config=ocr_config) if pytesseract else ''
                    except Exception:
                        ocr_text = ''
                    blocks.append(
                        f"\n{'='*80}\n"
                        f"[VISUAL ELEMENT #{count}]\n"
                        f"Location: Page {pno}, Image {idx}\n"
                        f"Type: {ext.upper()} image\n"
                        f"Dimensions: {w}x{h} pixels\n"
                        f"{'='*80}\n"
                        f"[OCR EXTRACTED TEXT FROM THIS VISUAL]:\n"
                        f"{ocr_text if ocr_text and ocr_text.strip() else '[No readable text detected in this image]'}\n"
                        f"{'='*80}\n"
                    )
                except Exception:
                    blocks.append(f"\n[VISUAL ELEMENT - ERROR]\nLocation: Page {pno}, Image {idx}\nError: See logs\n")
        doc.close()
        if count:
            return '\n'.join(blocks), count
        return '[No visual elements found in PDF]', 0
    except Exception:
        return '[Error extracting visual elements: see logs]', 0


def extract_text_from_pptx(pptx_path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(pptx_path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                try:
                    if hasattr(shape, 'text') and shape.text and shape.text.strip():
                        texts.append(shape.text)
                except Exception:
                    continue
            if texts:
                parts.append(f"\n{'='*80}\nSlide {i}\n{'='*80}\n" + '\n'.join(texts))
        return '\n'.join(parts)
    except Exception:
        try:
            with zipfile.ZipFile(pptx_path, 'r') as z:
                slide_files = [n for n in z.namelist() if n.startswith('ppt/slides/slide') and n.endswith('.xml')]
                slide_files.sort()
                out = []
                for idx, s in enumerate(slide_files, 1):
                    try:
                        xml = z.read(s)
                        root = ET.fromstring(xml)
                        texts = [t.text for t in root.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}t') if t.text and t.text.strip()]
                        if texts:
                            out.append(f"\n{'='*80}\nSlide {idx}\n{'='*80}\n" + '\n'.join(texts))
                    except Exception:
                        LOG.exception('Error parsing slide %s in %s', s, pptx_path)
                return '\n'.join(out)
        except Exception:
            LOG.exception('Error extracting text from pptx %s', pptx_path)
            return ''


def extract_images_from_pptx(pptx_path: Path, ocr_lang: Optional[str] = None, ocr_config: Optional[str] = None) -> Tuple[str, int]:
    try:
        with zipfile.ZipFile(pptx_path, 'r') as z:
            media = [n for n in z.namelist() if n.startswith('ppt/media/')]
            texts = []
            count = 0
            for m in media:
                try:
                    data = z.read(m)
                    image = Image.open(io.BytesIO(data)).convert('RGB')
                    count += 1
                    w, h = image.size
                    try:
                        ocr_text = pytesseract.image_to_string(image, lang=ocr_lang, config=ocr_config) if pytesseract else ''
                    except Exception:
                        ocr_text = ''
                    texts.append(
                        f"\n{'='*80}\n"
                        f"[VISUAL ELEMENT #{count}]\n"
                        f"Location: {m}\n"
                        f"Type: {Path(m).suffix.lstrip('.').upper()} image\n"
                        f"Dimensions: {w}x{h} pixels\n"
                        f"{'='*80}\n"
                        f"[OCR EXTRACTED TEXT FROM THIS VISUAL]:\n"
                        f"{ocr_text if ocr_text and ocr_text.strip() else '[No readable text detected in this image]'}\n"
                        f"{'='*80}\n"
                    )
                except Exception:
                    texts.append(f"\n[VISUAL ELEMENT - ERROR]\nLocation: {m}\nError: See logs\n")
            if count:
                return '\n'.join(texts), count
            return '[No visual elements found in PPTX]', 0
    except Exception:
        LOG.exception('Error extracting visual elements from %s', pptx_path)
        return '[Error extracting visual elements: see logs]', 0


def write_output_file(path: Path, metadata: str, text_section: str, visuals_section: str):
    content = (
        f"{'='*80}\n"
        f"[DOCUMENT METADATA]\n"
        f"{'='*80}\n"
        f"{metadata}\n"
        f"{'='*80}\n\n"
        f"{'#'*80}\n[SECTION 1: TEXT CONTENT]\n{'#'*80}\n"
        f"{text_section}\n\n"
        f"{'#'*80}\n[SECTION 2: VISUAL ELEMENTS - OCR EXTRACTED TEXT]\n{'#'*80}\n"
        f"{visuals_section}\n\n"
        f"{'='*80}\n[END OF DOCUMENT]\n{'='*80}\n"
    )
    path.write_text(content, encoding='utf-8')


def process_file(path: Path, output_dir: Path, ocr_lang: Optional[str] = None, ocr_config: Optional[str] = None):
    if path.suffix.lower() == '.pdf':
        LOG.info('Processing PDF: %s', path.name)
        text = extract_text_from_pdf(path)
        visuals, count = extract_images_from_pdf(path, ocr_lang=ocr_lang, ocr_config=ocr_config)
        metadata = f"FILE: {path.name}\nEXTRACTION DATE: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\nTYPE: PDF"
        out_path = output_dir / f"{path.stem}_extracted.txt"
        write_output_file(out_path, metadata, text, f"Total Visual Elements: {count}\n{visuals}")
        LOG.info('Wrote output: %s', out_path)
        return out_path
    elif path.suffix.lower() == '.pptx':
        LOG.info('Processing PPTX: %s', path.name)
        text = extract_text_from_pptx(path)
        visuals, count = extract_images_from_pptx(path, ocr_lang=ocr_lang, ocr_config=ocr_config)
        metadata = f"FILE: {path.name}\nEXTRACTION DATE: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\nTYPE: PPTX"
        out_path = output_dir / f"{path.stem}_extracted.txt"
        write_output_file(out_path, metadata, text, f"Total Visual Elements: {count}\n{visuals}")
        LOG.info('Wrote output: %s', out_path)
        return out_path
    else:
        LOG.debug('Unsupported file type: %s', path)
        return None


def main():
    parser = argparse.ArgumentParser(description='Extractor for Aston Docs (PDF & PPTX)')
    parser.add_argument('--dir', '-d', help='Directory to scan (default: parent of script)')
    parser.add_argument('--recursive', '-r', action='store_true', help='Recursively scan subdirectories')
    parser.add_argument('--ocr-lang', default=None, help='OCR language for pytesseract (e.g., eng)')
    parser.add_argument('--tesseract-cmd', default=None, help='Explicit tesseract executable path')
    parser.add_argument('--log-level', default='INFO', help='Logging level')
    args = parser.parse_args()

    script_dir = Path(__file__).parent
    if args.dir:
        base = Path(args.dir)
    else:
        base = script_dir
        for p in script_dir.parents:
            if p.name.lower() in ('aston docs', 'aston_docs'):
                base = p
                break

    recursive = args.recursive or (base.name.lower() in ('aston docs', 'aston_docs'))
    logging.info('Scanning: %s (recursive=%s)', base, recursive)

    if args.tesseract_cmd and pytesseract:
        pytesseract.pytesseract.tesseract_cmd = args.tesseract_cmd
        logging.info('Using tesseract from: %s', args.tesseract_cmd)

    output_dir = setup_output_directory(base)

    if recursive:
        files = list(base.rglob('*.pdf')) + list(base.rglob('*.pptx'))
    else:
        files = list(base.glob('*.pdf')) + list(base.glob('*.pptx'))

    if not files:
        logging.info('No PDF or PPTX files found in %s', base)
        return

    logging.info('Found %d files to process', len(files))

    processed = 0
    for f in files:
        try:
            res = process_file(f, output_dir, ocr_lang=args.ocr_lang, ocr_config='--psm 3')
            if res:
                processed += 1
        except Exception:
            logging.exception('Error processing %s', f)

    LOG.info('Done. Processed %d files. Outputs in %s', processed, output_dir)


if __name__ == '__main__':
    main()
