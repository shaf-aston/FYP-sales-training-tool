"""
Extracts text from PDF, DOCX, DOTX, PPTX files into extracted_texts/ folder.
Outputs: extracted_texts/ with project_management/ subfolder for SPM lectures.
"""
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
import shutil

try:
    import fitz
    PDF_OK = True
except ImportError:
    PDF_OK = False
    print("⚠️  pip install pymupdf")

try:
    from pptx import Presentation
    PPTX_OK = True
except ImportError:
    PPTX_OK = False
    print("⚠️  pip install python-pptx")

BASE = Path(r"c:\Users\Shaf\Downloads\Final Year Project pack folder\Sales roleplay chatbot\Aston Docs")
OUTPUT_DIR = BASE / "extracted_texts"
PM_OUTPUT_DIR = OUTPUT_DIR / "project_management"
TUTORIALS_OUTPUT_DIR = PM_OUTPUT_DIR / "tutorials"
WORD_NS = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'


def clean_filename(name):
    name = Path(name).stem.replace(" ", "_").replace("-", "_").replace("(", "").replace(")", "")
    while "__" in name:
        name = name.replace("__", "_")
    return name + ".txt"


def extract_pdf(file_path):
    if not PDF_OK:
        return "[ERROR: pymupdf not installed]"
    try:
        doc = fitz.open(file_path)
        text = [f"{'='*60}\nPAGE {i}\n{'='*60}\n{page.get_text()}" 
                for i, page in enumerate(doc, 1) if page.get_text().strip()]
        doc.close()
        return "\n\n".join(text)
    except Exception as e:
        return f"[ERROR: {e}]"


def extract_docx(file_path):
    """DOCX/DOTX are ZIP archives - extract text from word/document.xml"""
    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            if 'word/document.xml' not in zf.namelist():
                return ""
            root = ET.fromstring(zf.read('word/document.xml'))
            text = [''.join(t.text for t in para.iter(f'{WORD_NS}t') if t.text)
                    for para in root.iter(f'{WORD_NS}p')]
            return "\n\n".join(t for t in text if t.strip())
    except Exception as e:
        return f"[ERROR: {e}]"


def extract_pptx(file_path):
    if not PPTX_OK:
        return "[ERROR: python-pptx not installed]"
    try:
        prs = Presentation(file_path)
        text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = [f"{'='*60}\nSLIDE {i}\n{'='*60}"]
            slide_text.extend(shape.text for shape in slide.shapes 
                            if hasattr(shape, "text") and shape.text.strip())
            if len(slide_text) > 1:
                text.append("\n".join(slide_text))
        return "\n\n".join(text)
    except Exception as e:
        return f"[ERROR: {e}]"


def process_file(file_path, output_dir):
    file_path = Path(file_path)
    ext = file_path.suffix.lower()
    
    extractors = {'.pdf': extract_pdf, '.docx': extract_docx, 
                  '.dotx': extract_docx, '.pptx': extract_pptx}
    if ext not in extractors:
        return False
    
    content = extractors[ext](file_path)
    output_name = clean_filename(file_path.name)
    header = f"{'#'*70}\n# SOURCE: {file_path.name}\n# SIZE: {file_path.stat().st_size / 1024:.1f} KB\n{'#'*70}\n\n"
    
    (output_dir / output_name).write_text(header + content, encoding='utf-8')
    print(f"  ✓ {file_path.name} → {output_name}")
    return True


def main():
    print("=" * 70 + "\nUNIFIED DOCUMENT EXTRACTOR\n" + "=" * 70)
    
    if OUTPUT_DIR.exists():
        shutil.rmtree(OUTPUT_DIR)
    for d in [OUTPUT_DIR, PM_OUTPUT_DIR, TUTORIALS_OUTPUT_DIR]:
        d.mkdir(parents=True, exist_ok=True)
    
    print(f"\n📁 Output: {OUTPUT_DIR}\n\n📄 MAIN DOCUMENTS:")
    for f in sorted(BASE.iterdir()):
        if f.is_file():
            process_file(f, OUTPUT_DIR)
    
    pm_folder = BASE / "Project management"
    if pm_folder.exists():
        print("\n📄 PROJECT MANAGEMENT:")
        for f in sorted(pm_folder.iterdir()):
            if f.is_file():
                process_file(f, PM_OUTPUT_DIR)
        
        tutorials = pm_folder / "tutorials"
        if tutorials.exists():
            print("\n📄 TUTORIALS:")
            for f in sorted(tutorials.iterdir()):
                if f.is_file():
                    process_file(f, TUTORIALS_OUTPUT_DIR)
    
    total = sum(1 for _ in OUTPUT_DIR.rglob("*.txt"))
    total_size = sum(f.stat().st_size for f in OUTPUT_DIR.rglob("*.txt")) / 1024
    print(f"\n{'='*70}\n✅ Extracted {total} files ({total_size:.1f} KB) → {OUTPUT_DIR}")


if __name__ == "__main__":
    main()
